#!/usr/bin/env python3
"""
94CramManageSystem - 行銷推銷簡報生成器
莫蘭迪色系 + 現代風格 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ====== 莫蘭迪色系配色 ======
COLORS = {
    'primary':      RGBColor(0x4A, 0x6B, 0x8A),   # 深莫蘭迪藍
    'secondary':    RGBColor(0x8B, 0x9D, 0x83),   # 莫蘭迪綠
    'accent':       RGBColor(0xC4, 0x8B, 0x6A),   # 莫蘭迪橘
    'accent2':      RGBColor(0xA0, 0x7E, 0x93),   # 莫蘭迪紫
    'dark':         RGBColor(0x2D, 0x3A, 0x4A),   # 深色背景
    'dark2':        RGBColor(0x3A, 0x4A, 0x5C),   # 次深色
    'light':        RGBColor(0xF5, 0xF0, 0xEB),   # 淺米色
    'light2':       RGBColor(0xE8, 0xE0, 0xD8),   # 次淺色
    'white':        RGBColor(0xFF, 0xFF, 0xFF),
    'text_dark':    RGBColor(0x2D, 0x2D, 0x2D),
    'text_light':   RGBColor(0x6B, 0x6B, 0x6B),
    'red':          RGBColor(0xC0, 0x5C, 0x5C),   # 莫蘭迪紅
    'gold':         RGBColor(0xC4, 0xA3, 0x5A),   # 莫蘭迪金
    'green_check':  RGBColor(0x5A, 0x8C, 0x6A),   # 打勾綠
    'red_cross':    RGBColor(0xB0, 0x5A, 0x5A),   # 叉叉紅
    'gradient_top': RGBColor(0x2D, 0x3A, 0x4A),
    'gradient_bot': RGBColor(0x4A, 0x6B, 0x8A),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====== 工具函數 ======
def add_bg(slide, color):
    """設定整頁背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    """加入矩形色塊"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha_el = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', str(int(alpha * 1000)))
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    """加入圓角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=COLORS['text_dark'],
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft JhengHei'):
    """加入文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=16, color=COLORS['text_dark'], bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name='Microsoft JhengHei'):
    """在既有 text_frame 加入段落"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_circle(slide, left, top, size, color):
    """加入圓形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, bg_color, icon_color):
    """加入帶圖標的卡片"""
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    # 圖標圓形
    circle = add_circle(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.7), icon_color)
    # 圖標文字
    add_text(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.7), Inches(0.7),
             icon_text, font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    # 標題
    add_text(slide, left + Inches(1.15), top + Inches(0.3), width - Inches(1.5), Inches(0.5),
             title, font_size=18, color=COLORS['dark'], bold=True)
    # 描述
    add_text(slide, left + Inches(1.15), top + Inches(0.75), width - Inches(1.5), height - Inches(1.0),
             desc, font_size=13, color=COLORS['text_light'])
    return card

def add_stat_card(slide, left, top, number, label, color):
    """加入數據統計卡片"""
    card = add_rounded_rect(slide, left, top, Inches(2.4), Inches(1.6), COLORS['white'])
    # 頂部色條
    add_rect(slide, left, top, Inches(2.4), Inches(0.06), color)
    # 數字
    add_text(slide, left, top + Inches(0.25), Inches(2.4), Inches(0.8),
             number, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 標籤
    add_text(slide, left, top + Inches(1.0), Inches(2.4), Inches(0.5),
             label, font_size=14, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)
    return card

def add_feature_bullet(text_frame, icon, text, font_size=15, color=COLORS['text_dark']):
    """加入功能要點"""
    p = text_frame.add_paragraph()
    p.text = f"{icon}  {text}"
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = 'Microsoft JhengHei'
    p.space_before = Pt(6)
    p.space_after = Pt(2)
    return p

def slide_header(slide, title, subtitle=None):
    """統一頁面標題"""
    # 頂部裝飾條
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLORS['primary'])
    # 標題
    add_text(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             title, font_size=32, color=COLORS['dark'], bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
                 subtitle, font_size=16, color=COLORS['text_light'])
    # 右上角品牌
    add_text(slide, Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.4),
             '94Cram 智慧補教', font_size=14, color=COLORS['primary'], bold=True, alignment=PP_ALIGN.RIGHT)

# =========================================================
# SLIDE 1: 封面
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['dark'])

# 裝飾元素
add_circle(slide, Inches(-1.5), Inches(-2), Inches(6), COLORS['dark2'])
add_circle(slide, Inches(9), Inches(4), Inches(5), COLORS['dark2'])

# 頂部色條
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), COLORS['accent'])

# 品牌標識
add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
         '94Cram', font_size=24, color=COLORS['accent'], bold=True)

# 主標題
add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         '智慧補教管理生態系', font_size=54, color=COLORS['white'], bold=True)

# 副標題
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
         '三大系統 × AI 驅動 × 一站式解決方案', font_size=28, color=COLORS['light2'])

# 分隔線
add_rect(slide, Inches(1), Inches(4.3), Inches(3), Inches(0.04), COLORS['accent'])

# 描述
add_text(slide, Inches(1), Inches(4.6), Inches(8), Inches(0.5),
         '學員管理 ｜ 智慧點名 ｜ 庫存管控 ｜ AI 助手 ｜ LINE/Telegram Bot',
         font_size=18, color=COLORS['light2'])

# 日期
add_text(slide, Inches(1), Inches(5.5), Inches(5), Inches(0.4),
         '2026 產品介紹', font_size=16, color=COLORS['text_light'])

# 右下角裝飾
add_rounded_rect(slide, Inches(9.5), Inches(5.5), Inches(3), Inches(1.3), COLORS['primary'])
add_text(slide, Inches(9.5), Inches(5.7), Inches(3), Inches(0.4),
         '免費試用 30 天', font_size=20, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.5), Inches(6.15), Inches(3), Inches(0.3),
         '零硬體投資 · 即開即用', font_size=14, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 2: 補習班的痛點
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '補習班經營的六大痛點', '您是否正在面對這些挑戰？')

pain_points = [
    ('😰', '學員管理混亂', '紙本名冊、Excel 表格散落各處\n學員資料不統一，查詢耗時', COLORS['red']),
    ('📋', '點名效率低落', '每堂課手動點名浪費 5-10 分鐘\n家長無法即時知道孩子出席狀況', COLORS['accent']),
    ('💸', '收費漏洞百出', '繳費記錄靠人工，漏收時有所聞\n催款困難，應收帳款追蹤困難', COLORS['accent2']),
    ('📦', '教材管理失控', '講義庫存靠感覺，常缺貨或囤積\n各校區教材流向不透明', COLORS['primary']),
    ('🚪', '學員流失無感', '學員流失才發現為時已晚\n缺乏預警機制，流失原因難追蹤', COLORS['red']),
    ('🔐', '資料安全堪憂', '重要資料存在本機硬碟\n沒有備份，電腦一壞全部歸零', COLORS['dark2']),
]

for i, (icon, title, desc, color) in enumerate(pain_points):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.7 + row * 2.7)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.3), COLORS['white'])
    # 左側色條
    add_rect(slide, left, top + Inches(0.3), Inches(0.06), Inches(1.7), color)

    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(0.6), Inches(0.6),
             icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.9), top + Inches(0.25), Inches(2.7), Inches(0.4),
             title, font_size=18, color=color, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.85), Inches(3.2), Inches(1.2),
             desc, font_size=13, color=COLORS['text_light'])


# =========================================================
# SLIDE 3: 解決方案總覽
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '94Cram 一站式解決方案', '三大系統 + AI 助手，完整覆蓋補習班營運需求')

systems = [
    ('管', '94Manage', '學員管理系統', '學員資料 · 課程管理 · 招生漏斗\n收費帳務 · 薪資管理 · AI 流失預警\n成績追蹤 · 知識庫 · 報表分析', COLORS['primary']),
    ('名', '94inClass', '智慧點名系統', 'NFC 刷卡點名（1秒完成）\nAI 臉部辨識 · 即時家長通知\n出勤統計 · 請假管理 · 成績管理\n課表排程 · 繳費管理', COLORS['secondary']),
    ('庫', '94Stock', '庫存管理系統', '多倉庫管理 · 條碼掃描\n進出貨追蹤 · 低庫存預警\n採購訂單流程 · 供應商管理\n盤點作業 · AI 備貨預測', COLORS['accent']),
]

for i, (icon_char, name, subtitle, features, color) in enumerate(systems):
    left = Inches(0.5 + i * 4.2)
    top = Inches(1.7)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(4.5), COLORS['white'])
    # 頂部色帶
    add_rect(slide, left, top, Inches(3.8), Inches(0.8), color)
    # 圖標
    circle = add_circle(slide, left + Inches(1.4), top + Inches(0.08), Inches(0.65), COLORS['white'])
    add_text(slide, left + Inches(1.4), top + Inches(0.08), Inches(0.65), Inches(0.65),
             icon_char, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 系統名
    add_text(slide, left, top + Inches(0.95), Inches(3.8), Inches(0.5),
             name, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 副標題
    add_text(slide, left, top + Inches(1.4), Inches(3.8), Inches(0.4),
             subtitle, font_size=14, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)
    # 分隔線
    add_rect(slide, left + Inches(0.5), top + Inches(1.85), Inches(2.8), Inches(0.02), COLORS['light2'])
    # 功能列表
    add_text(slide, left + Inches(0.4), top + Inches(2.0), Inches(3.0), Inches(2.3),
             features, font_size=13, color=COLORS['text_dark'])

# AI 底部橫幅
add_rounded_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.85), COLORS['dark'])
add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.35),
         '🤖  AI 驅動核心：頂尖大型語言模型 + 智慧知識引擎 + 自然語言操作 + 智慧預測',
         font_size=17, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.85), Inches(10), Inches(0.3),
         'Telegram / LINE 聊天即操作，家長學員零門檻使用',
         font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 4: 94Manage 學員管理 — 核心功能
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '94Manage — 學員管理系統', '從招生到畢業，全生命週期管理')

# 左半部分：招生漏斗
add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.5), COLORS['white'])
add_text(slide, Inches(0.8), Inches(1.75), Inches(5), Inches(0.4),
         '🎯 招生漏斗 & 智慧分析', font_size=20, color=COLORS['primary'], bold=True)

funnel_stages = [
    ('新諮詢', '100%', Inches(5.0), COLORS['primary']),
    ('已聯絡', '75%', Inches(4.3), COLORS['secondary']),
    ('預約試聽', '50%', Inches(3.6), COLORS['accent']),
    ('完成試聽', '35%', Inches(2.9), COLORS['accent2']),
    ('正式報名', '25%', Inches(2.2), COLORS['green_check']),
]

for i, (stage, pct, width, color) in enumerate(funnel_stages):
    t = Inches(2.4 + i * 0.75)
    offset = (Inches(5.0) - width) / 2
    left = Inches(1.0) + offset
    add_rounded_rect(slide, left, t, width, Inches(0.55), color)
    add_text(slide, left, t + Inches(0.05), width, Inches(0.45),
             f'{stage}  {pct}', font_size=14, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.5),
         '自動追蹤每階段轉換率 · 顧問績效排名 · 預期營收計算',
         font_size=12, color=COLORS['text_light'])

# 右半部分：核心模組
right_features = [
    ('👨‍🎓', '學員管理', '資料建檔 · 狀態追蹤 · 批量匯入'),
    ('📚', '課程管理', '五種收費模式 · 動態費率調整'),
    ('💰', '帳務系統', '繳費追蹤 · AI 自動發票 · 逾期提醒'),
    ('💼', '薪資管理', '自動計算 · 時薪/獎金 · 薪資單'),
    ('📊', '成績分析', '成績登錄 · 趨勢圖表 · 進步追蹤'),
    ('📄', 'AI 報表', '分校報告自動生成 · 學員詳細報告'),
]

for i, (icon, title, desc) in enumerate(right_features):
    top = Inches(1.65 + i * 0.9)
    card = add_rounded_rect(slide, Inches(6.8), top, Inches(5.8), Inches(0.78), COLORS['white'])
    add_text(slide, Inches(7.0), top + Inches(0.05), Inches(0.5), Inches(0.5),
             icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.6), top + Inches(0.05), Inches(1.5), Inches(0.35),
             title, font_size=16, color=COLORS['dark'], bold=True)
    add_text(slide, Inches(7.6), top + Inches(0.38), Inches(4.8), Inches(0.35),
             desc, font_size=12, color=COLORS['text_light'])


# =========================================================
# SLIDE 5: AI 流失預警系統
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '🧠 AI 流失預警系統', '業界首創：在學員流失前主動預警')

# 左邊說明
add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5), COLORS['white'])
add_text(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.4),
         '多維度智慧分析', font_size=20, color=COLORS['primary'], bold=True)

signals = [
    ('📉', '出勤率下降', '連續缺課或出勤率低於班級平均'),
    ('💳', '繳費延遲', '逾期繳費或拖延天數增加'),
    ('📊', '成績下降', '連續退步或大幅落後'),
    ('🔕', '互動減少', '課堂參與度降低'),
    ('📅', '請假頻繁', '請假次數異常增加'),
]

for i, (icon, title, desc) in enumerate(signals):
    top = Inches(2.4 + i * 0.8)
    add_text(slide, Inches(0.8), top, Inches(0.5), Inches(0.4), icon, font_size=20)
    add_text(slide, Inches(1.4), top, Inches(2), Inches(0.35),
             title, font_size=15, color=COLORS['dark'], bold=True)
    add_text(slide, Inches(1.4), top + Inches(0.3), Inches(4.5), Inches(0.35),
             desc, font_size=12, color=COLORS['text_light'])

# 右邊風險儀表板
add_rounded_rect(slide, Inches(6.6), Inches(1.6), Inches(6.2), Inches(5.5), COLORS['dark'])
add_text(slide, Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.4),
         '風險儀表板', font_size=20, color=COLORS['white'], bold=True)

risk_levels = [
    ('🔴', '高風險', '3 名學員', '立即聯繫 · 安排面談 · 提供優惠', COLORS['red']),
    ('🟡', '中風險', '8 名學員', '加強關懷 · 追蹤狀態 · 觀察趨勢', COLORS['gold']),
    ('🟢', '低風險', '45 名學員', '維持現狀 · 定期關懷', COLORS['green_check']),
]

for i, (icon, level, count, action, color) in enumerate(risk_levels):
    top = Inches(2.5 + i * 1.5)
    add_rounded_rect(slide, Inches(7.0), top, Inches(5.4), Inches(1.2), COLORS['dark2'])
    add_rect(slide, Inches(7.0), top, Inches(0.08), Inches(1.2), color)
    add_text(slide, Inches(7.3), top + Inches(0.1), Inches(0.4), Inches(0.4),
             icon, font_size=22)
    add_text(slide, Inches(7.8), top + Inches(0.1), Inches(1.5), Inches(0.35),
             level, font_size=18, color=color, bold=True)
    add_text(slide, Inches(10.0), top + Inches(0.15), Inches(2), Inches(0.3),
             count, font_size=16, color=COLORS['white'], alignment=PP_ALIGN.RIGHT)
    add_text(slide, Inches(7.3), top + Inches(0.6), Inches(5), Inches(0.4),
             f'建議行動：{action}', font_size=12, color=COLORS['light2'])

# 底部統計
add_text(slide, Inches(7.0), Inches(6.4), Inches(5.5), Inches(0.35),
         '⚡ 每日自動掃描 · 準確率 92% · 提前 2-4 週預警',
         font_size=14, color=COLORS['gold'], bold=True, alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 6: 94inClass 點名系統
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '94inClass — 智慧點名系統', '1 秒完成點名，家長即時收到通知')

# NFC 點名卡片
add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(3.2), COLORS['primary'])
add_text(slide, Inches(0.5), Inches(1.85), Inches(3.8), Inches(0.5),
         '📱', font_size=40, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.5), Inches(3.8), Inches(0.5),
         'NFC 感應點名', font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(3.1), Inches(3.2), Inches(0.35),
         '刷卡即到 · 1 秒完成', font_size=16, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(3.55), Inches(3.2), Inches(0.8),
         '市售 NFC 讀卡機 NT$300 即可\n支援所有 NFC 卡片\n零學習成本', font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

# AI 臉辨卡片
add_rounded_rect(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(3.2), COLORS['secondary'])
add_text(slide, Inches(4.7), Inches(1.85), Inches(3.8), Inches(0.5),
         '🤖', font_size=40, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.7), Inches(2.5), Inches(3.8), Inches(0.5),
         'AI 臉部辨識', font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(3.1), Inches(3.2), Inches(0.35),
         '走進教室自動辨識', font_size=16, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(3.55), Inches(3.2), Inches(0.8),
         '使用一般網路攝影機\n防代簽驗證\n科技感十足', font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

# LINE 通知卡片
add_rounded_rect(slide, Inches(8.9), Inches(1.6), Inches(3.8), Inches(3.2), COLORS['accent'])
add_text(slide, Inches(8.9), Inches(1.85), Inches(3.8), Inches(0.5),
         '💬', font_size=40, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.9), Inches(2.5), Inches(3.8), Inches(0.5),
         '即時家長通知', font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.2), Inches(3.1), Inches(3.2), Inches(0.35),
         'LINE 推播零延遲', font_size=16, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(9.2), Inches(3.55), Inches(3.2), Inches(0.8),
         '到校即通知家長\n遲到/缺席自動推播\n每日出勤摘要', font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

# 底部功能列
features_bottom = [
    ('📊 出勤統計', '日/週/月報表\n出勤率分析'),
    ('📝 請假管理', '線上請假申請\n額度追蹤'),
    ('📅 課表管理', '教室 · 老師 · 時段\n排程管理'),
    ('💳 繳費管理', '多元週期 · 狀態追蹤\n批次登錄'),
    ('📄 成績管理', '成績登錄 · 統計\n排名 · 進步追蹤'),
]

for i, (title, desc) in enumerate(features_bottom):
    left = Inches(0.5 + i * 2.55)
    top = Inches(5.2)
    card = add_rounded_rect(slide, left, top, Inches(2.3), Inches(1.8), COLORS['white'])
    add_text(slide, left, top + Inches(0.15), Inches(2.3), Inches(0.35),
             title, font_size=14, color=COLORS['dark'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), top + Inches(0.55), Inches(2.0), Inches(1.0),
             desc, font_size=12, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 7: 94Stock 庫存管理
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '94Stock — 庫存管理系統', '教材管理數位化，再也不怕講義缺貨')

# 主要流程
flow_items = [
    ('📦', '進貨', '採購收貨\n條碼掃描\n自動入帳'),
    ('➡️', '', ''),
    ('🏢', '倉儲', '多倉庫管理\n即時庫存\n安全存量'),
    ('➡️', '', ''),
    ('📤', '出貨', '班級領用\n銷售出貨\n簽收記錄'),
    ('➡️', '', ''),
    ('🔄', '轉庫', '校區間轉移\n雙邊自動帳\n交易記錄'),
]

x_pos = Inches(0.3)
for i, (icon, title, desc) in enumerate(flow_items):
    if title == '':
        # 箭頭
        add_text(slide, x_pos, Inches(2.5), Inches(0.6), Inches(0.6),
                 '→', font_size=36, color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.6)
    else:
        card = add_rounded_rect(slide, x_pos, Inches(1.6), Inches(2.7), Inches(2.6), COLORS['white'])
        add_text(slide, x_pos, Inches(1.75), Inches(2.7), Inches(0.5),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
        add_text(slide, x_pos, Inches(2.3), Inches(2.7), Inches(0.4),
                 title, font_size=20, color=COLORS['dark'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x_pos + Inches(0.2), Inches(2.8), Inches(2.3), Inches(1.2),
                 desc, font_size=13, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.9)

# 底部特色功能
bottom_features = [
    ('🔔', '低庫存預警', 'Telegram 即時推播\n低於安全存量自動通知', COLORS['red']),
    ('📋', '採購訂單流程', '草稿→審核→核准→收貨\n完整審批流程', COLORS['primary']),
    ('📊', '盤點管理', '建立盤點單 · 掃描盤點\n差異自動調整', COLORS['secondary']),
    ('🤖', 'AI 備貨預測', '基於歷史數據\n自動建議補貨量', COLORS['accent']),
]

for i, (icon, title, desc, color) in enumerate(bottom_features):
    left = Inches(0.5 + i * 3.2)
    top = Inches(4.7)
    card = add_rounded_rect(slide, left, top, Inches(2.9), Inches(2.4), COLORS['white'])
    add_rect(slide, left, top, Inches(2.9), Inches(0.06), color)
    add_text(slide, left, top + Inches(0.2), Inches(2.9), Inches(0.5),
             icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.75), Inches(2.9), Inches(0.4),
             title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(1.2), Inches(2.5), Inches(1.0),
             desc, font_size=12, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 8: AI 機器人 — Telegram & LINE
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['dark'])

# 頂部裝飾
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLORS['accent'])
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         '🤖 AI 智慧機器人', font_size=32, color=COLORS['white'], bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         '用「說話」管理補習班 — Telegram / LINE 自然語言操作', font_size=16, color=COLORS['light2'])
add_text(slide, Inches(10.5), Inches(0.45), Inches(2.5), Inches(0.4),
         '94Cram 智慧補教', font_size=14, color=COLORS['primary'], bold=True, alignment=PP_ALIGN.RIGHT)

# 管理員模式
add_rounded_rect(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.3), COLORS['dark2'])
add_text(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.4),
         '👔 管理員模式（Telegram）', font_size=18, color=COLORS['accent'], bold=True)

admin_commands = [
    ('「陳小明請假」', '→ 自動記錄請假申請'),
    ('「高二班繳 5000 元」', '→ 自動建立繳費紀錄'),
    ('「203 號講義剩幾本」', '→ 即時查詢庫存數量'),
    ('「今天出勤率多少」', '→ 顯示當日出勤統計'),
    ('「本月營收報表」', '→ 產生收入分析報告'),
    ('「切換到中壢分校」', '→ 切換操作分校'),
]

for i, (cmd, result) in enumerate(admin_commands):
    top = Inches(2.5 + i * 0.68)
    # 指令氣泡
    add_rounded_rect(slide, Inches(0.8), top, Inches(2.5), Inches(0.5), COLORS['primary'])
    add_text(slide, Inches(0.9), top + Inches(0.05), Inches(2.3), Inches(0.4),
             cmd, font_size=13, color=COLORS['white'], bold=True)
    # 回應
    add_text(slide, Inches(3.5), top + Inches(0.05), Inches(3.0), Inches(0.4),
             result, font_size=13, color=COLORS['light2'])

# 安全機制
add_rounded_rect(slide, Inches(0.8), Inches(6.4), Inches(5.5), Inches(0.45), COLORS['accent'])
add_text(slide, Inches(1.0), Inches(6.45), Inches(5.0), Inches(0.35),
         '🔒  寫入操作二次確認 · 防誤操作設計', font_size=13, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

# 家長模式
add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.3), COLORS['dark2'])
add_text(slide, Inches(7.1), Inches(1.9), Inches(5.5), Inches(0.4),
         '👨‍👩‍👧 家長模式（Telegram / LINE）', font_size=18, color=COLORS['secondary'], bold=True)

parent_features = [
    ('綁定驗證', '孩子姓名 + 電話末 4 碼\n簡單安全的身份驗證'),
    ('出勤查詢', '「小明今天有到嗎？」\n即時回覆到校狀態'),
    ('成績查詢', '「小明最近考試成績」\n自動顯示成績與排名'),
    ('費用查詢', '「小明學費繳清了嗎？」\n顯示繳費狀態與明細'),
    ('課表查詢', '「小明這週上課時間」\n完整課表一目瞭然'),
]

for i, (title, desc) in enumerate(parent_features):
    top = Inches(2.5 + i * 0.9)
    add_rounded_rect(slide, Inches(7.1), top, Inches(5.4), Inches(0.75), COLORS['dark'])
    add_text(slide, Inches(7.3), top + Inches(0.05), Inches(1.5), Inches(0.3),
             title, font_size=14, color=COLORS['secondary'], bold=True)
    add_text(slide, Inches(7.3), top + Inches(0.32), Inches(5.0), Inches(0.4),
             desc, font_size=11, color=COLORS['light2'])

# 家長底部
add_rounded_rect(slide, Inches(7.1), Inches(6.4), Inches(5.5), Inches(0.45), COLORS['secondary'])
add_text(slide, Inches(7.3), Inches(6.45), Inches(5.0), Inches(0.35),
         '💡  零學習成本 · 用對話就能查詢一切', font_size=13, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 9: 與競品比較
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '競品比較分析', '為什麼 94Cram 是最佳選擇？')

# 表頭
headers = ['功能比較', '94Cram\n智慧補教', '傳統補教\nERP', 'Excel\n人工管理', '其他 SaaS\n管理系統']
header_colors = [COLORS['dark'], COLORS['primary'], COLORS['text_light'], COLORS['text_light'], COLORS['text_light']]
col_widths = [Inches(3.0), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]

for i, (header, color) in enumerate(zip(headers, header_colors)):
    left = Inches(0.5) + sum(w for w in [Emu(0)] + list(col_widths[:i]))
    bg_color = COLORS['primary'] if i == 1 else COLORS['dark']
    add_rounded_rect(slide, left, Inches(1.6), col_widths[i], Inches(0.7), bg_color)
    add_text(slide, left, Inches(1.62), col_widths[i], Inches(0.65),
             header, font_size=13, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

# 比較項目
compare_items = [
    ('學員管理', '✓', '✓', '△', '✓'),
    ('智慧點名 (NFC/臉辨)', '✓', '✗', '✗', '△'),
    ('AI 流失預警', '✓', '✗', '✗', '✗'),
    ('庫存管理', '✓', '△', '✗', '✗'),
    ('LINE/Telegram Bot', '✓', '✗', '✗', '✗'),
    ('AI 自然語言操作', '✓', '✗', '✗', '✗'),
    ('家長即時通知', '✓', '△', '✗', '△'),
    ('多分校支援', '✓', '✓', '✗', '△'),
    ('雲端 SaaS（免硬體）', '✓', '✗', '—', '✓'),
    ('AI 備貨預測', '✓', '✗', '✗', '✗'),
    ('月費 (參考)', 'NT$2,999 起', 'NT$10,000+', '免費', 'NT$5,000+'),
    ('建置費', '0 元', '10~50 萬', '0 元', '0~5 萬'),
]

for row_i, (feature, *values) in enumerate(compare_items):
    top = Inches(2.35 + row_i * 0.4)
    bg = COLORS['white'] if row_i % 2 == 0 else COLORS['light2']

    # 功能名稱
    left = Inches(0.5)
    add_rect(slide, left, top, col_widths[0], Inches(0.38), bg)
    add_text(slide, left + Inches(0.2), top, col_widths[0], Inches(0.38),
             feature, font_size=12, color=COLORS['text_dark'], bold=True)

    # 數值列
    for col_i, val in enumerate(values):
        left = Inches(0.5) + sum(w for w in [Emu(0)] + list(col_widths[:col_i + 1]))
        cell_bg = RGBColor(0xEE, 0xF5, 0xF0) if col_i == 0 and val == '✓' else bg
        add_rect(slide, left, top, col_widths[col_i + 1], Inches(0.38), cell_bg)

        if val == '✓':
            c = COLORS['green_check']
            display = '✓'
        elif val == '✗':
            c = COLORS['red_cross']
            display = '✗'
        elif val == '△':
            c = COLORS['gold']
            display = '△'
        else:
            c = COLORS['text_dark']
            display = val

        font_bold = True if val in ('✓', '✗', '△') else False
        fs = 14 if val in ('✓', '✗', '△') else 11
        add_text(slide, left, top, col_widths[col_i + 1], Inches(0.38),
                 display, font_size=fs, color=c, bold=font_bold, alignment=PP_ALIGN.CENTER)

# 底部結論
add_rounded_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8), COLORS['primary'])
add_text(slide, Inches(1.0), Inches(6.48), Inches(11), Inches(0.3),
         '💡 94Cram 是市場上唯一整合「學員管理 + 點名 + 庫存 + AI + Bot」的補教管理系統',
         font_size=16, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(6.82), Inches(11), Inches(0.25),
         '其他系統至少需要 3-4 套軟體才能達到相同效果，且無 AI 智慧功能',
         font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 10: 技術架構優勢
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '技術架構優勢', '企業級雲端架構，補習班級的價格')

# 架構圖示
arch_layers = [
    ('使用者端', 'Web 瀏覽器 · LINE · Telegram · 手機', COLORS['primary']),
    ('前端層', '新一代響應式框架 · 強型別全棧開發 · 伺服器端渲染加速', COLORS['secondary']),
    ('API 層', '軍規級身份認證 · 角色權限管控 · 資料驗證防護 · 標準化介面', COLORS['accent']),
    ('AI 層', '最新一代大型語言模型 · 智慧知識引擎 · 語意向量檢索 · 意圖理解', COLORS['accent2']),
    ('資料層', '企業級關聯式資料庫 · 型別安全 ORM · 多租戶隔離 · 審計日誌', COLORS['dark2']),
    ('基礎設施', '頂級雲端無伺服器架構 · 託管式資料庫 · 自動擴縮 · 零停機部署', COLORS['primary']),
]

for i, (layer, desc, color) in enumerate(arch_layers):
    top = Inches(1.55 + i * 0.9)
    # 層級標籤
    add_rounded_rect(slide, Inches(0.5), top, Inches(2.0), Inches(0.7), color)
    add_text(slide, Inches(0.5), top + Inches(0.1), Inches(2.0), Inches(0.5),
             layer, font_size=15, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    # 描述橫條
    add_rounded_rect(slide, Inches(2.7), top, Inches(5.3), Inches(0.7), COLORS['white'])
    add_text(slide, Inches(2.9), top + Inches(0.1), Inches(5.0), Inches(0.5),
             desc, font_size=13, color=COLORS['text_dark'])

# 右邊優勢列表
advantages = [
    ('⚡', '更新零停機', '雙軌熱切換部署\n更新系統完全不影響使用'),
    ('🔒', '企業級安全', '軍規級認證 + 角色權限\n每筆操作可追蹤'),
    ('📈', '自動擴縮', '尖峰時段自動擴展\n離峰自動縮減省錢'),
    ('🌏', '台灣機房', '國際頂級雲端台灣區\n延遲 < 50ms 超快速'),
    ('💾', '自動備份', '每日自動備份\n資料永不遺失'),
    ('🔐', '資料隔離', '多租戶架構\n每家補習班資料獨立'),
]

for i, (icon, title, desc) in enumerate(advantages):
    col = i % 2
    row = i // 2
    left = Inches(8.3 + col * 2.5)
    top = Inches(1.55 + row * 1.85)

    card = add_rounded_rect(slide, left, top, Inches(2.3), Inches(1.65), COLORS['white'])
    add_text(slide, left, top + Inches(0.1), Inches(2.3), Inches(0.4),
             icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2.3), Inches(0.35),
             title, font_size=14, color=COLORS['dark'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.9), Inches(2.1), Inches(0.65),
             desc, font_size=11, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 11: 安全與合規
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['dark'])
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLORS['accent'])
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         '🔐 安全與合規', font_size=32, color=COLORS['white'], bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         '補教業最嚴謹的資安標準', font_size=16, color=COLORS['light2'])
add_text(slide, Inches(10.5), Inches(0.45), Inches(2.5), Inches(0.4),
         '94Cram 智慧補教', font_size=14, color=COLORS['primary'], bold=True, alignment=PP_ALIGN.RIGHT)

security_items = [
    ('🔑', '軍規級認證 + SSO', '三系統單一登入\n一組帳號通用全平台\n自動逾時登出'),
    ('👥', 'RBAC 角色權限', '6 種角色細粒度控管\n資源級存取控制\n每個人只看到該看的'),
    ('📝', '完整審計日誌', '所有操作留下紀錄\n何人何時改了什麼\nIP 來源追蹤'),
    ('🏠', '多租戶隔離', '每家補習班資料獨立\n互不干擾、不外洩\n嚴格 tenantId 驗證'),
    ('☁️', '頂級雲端防護', 'SSL/TLS 加密傳輸\n台灣機房資料主權\n定期自動備份'),
    ('🛡️', '個資保護', '符合個資法規範\n敏感資料加密存儲\n環境變數管理密鑰'),
]

for i, (icon, title, desc) in enumerate(security_items):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.7 + row * 2.7)

    card = add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.4), COLORS['dark2'])
    add_text(slide, left, top + Inches(0.15), Inches(3.8), Inches(0.5),
             icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.7), Inches(3.8), Inches(0.4),
             title, font_size=18, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.3), top + Inches(1.2), Inches(3.2), Inches(1.0),
             desc, font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 12: 成本效益
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '成本效益分析', '最低成本，最高效率')

# 數據卡片
stats = [
    ('NT$0', '建置費', COLORS['green_check']),
    ('NT$300', 'NFC 讀卡機', COLORS['primary']),
    ('30 天', '免費試用', COLORS['accent']),
    ('99.95%', '系統可用性', COLORS['secondary']),
]

for i, (num, label, color) in enumerate(stats):
    add_stat_card(slide, Inches(0.5 + i * 3.15), Inches(1.5), num, label, color)

# 成本比較表
add_rounded_rect(slide, Inches(0.5), Inches(3.5), Inches(6.0), Inches(3.7), COLORS['white'])
add_text(slide, Inches(0.8), Inches(3.65), Inches(5), Inches(0.4),
         '💰 3 年總成本比較（TCO）', font_size=18, color=COLORS['dark'], bold=True)

cost_compare = [
    ('項目', '94Cram', '傳統系統'),
    ('建置費', 'NT$0', 'NT$100,000+'),
    ('硬體採購', 'NT$300', 'NT$50,000+'),
    ('年度授權', 'NT$35,988', 'NT$120,000'),
    ('維護費用', 'NT$0', 'NT$30,000/年'),
    ('3 年總計', 'NT$108,264', 'NT$430,000+'),
    ('節省', '— —', '75%↓'),
]

for row_i, (item, ours, theirs) in enumerate(cost_compare):
    top = Inches(4.15 + row_i * 0.42)
    is_header = row_i == 0
    is_total = row_i >= 5
    bg = COLORS['primary'] if is_header else (RGBColor(0xEE, 0xF5, 0xF0) if is_total else (COLORS['white'] if row_i % 2 == 0 else COLORS['light']))
    fc = COLORS['white'] if is_header else (COLORS['green_check'] if is_total else COLORS['text_dark'])

    add_rect(slide, Inches(0.8), top, Inches(2.0), Inches(0.38), bg)
    add_text(slide, Inches(0.9), top, Inches(1.8), Inches(0.38),
             item, font_size=12, color=fc, bold=is_header or is_total)

    add_rect(slide, Inches(2.8), top, Inches(1.7), Inches(0.38), bg)
    add_text(slide, Inches(2.8), top, Inches(1.7), Inches(0.38),
             ours, font_size=12, color=fc, bold=is_header or is_total, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.5), top, Inches(1.7), Inches(0.38), bg)
    add_text(slide, Inches(4.5), top, Inches(1.7), Inches(0.38),
             theirs, font_size=12, color=COLORS['white'] if is_header else COLORS['red_cross'],
             bold=is_header, alignment=PP_ALIGN.CENTER)

# ROI 面板
add_rounded_rect(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.7), COLORS['primary'])
add_text(slide, Inches(7.2), Inches(3.7), Inches(5), Inches(0.4),
         '📈 導入效益', font_size=20, color=COLORS['white'], bold=True)

roi_items = [
    ('每堂省 5 分鐘', '智慧點名取代手動點名', '每月省 10+ 小時'),
    ('學員流失率降低', 'AI 預警提前 2-4 週通知', '留住 15% 潛在流失學員'),
    ('收費零遺漏', '自動追蹤每筆繳費', '每月多收 NT$5,000+'),
    ('教材零浪費', 'AI 備貨預測精準採購', '減少 20% 囤積浪費'),
    ('招生轉換率提升', '漏斗分析優化招生流程', '轉換率提升 30%'),
]

for i, (title, desc, result) in enumerate(roi_items):
    top = Inches(4.25 + i * 0.62)
    add_text(slide, Inches(7.2), top, Inches(2.0), Inches(0.3),
             title, font_size=13, color=COLORS['white'], bold=True)
    add_text(slide, Inches(7.2), top + Inches(0.25), Inches(3.0), Inches(0.25),
             desc, font_size=11, color=COLORS['light2'])
    add_text(slide, Inches(10.5), top + Inches(0.05), Inches(2.0), Inches(0.3),
             result, font_size=12, color=COLORS['gold'], bold=True, alignment=PP_ALIGN.RIGHT)


# =========================================================
# SLIDE 13: 服務方案
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '服務方案', '彈性選擇，隨需擴展')

plans = [
    ('入門版', 'NT$2,999/月', '50 學員以下', COLORS['secondary'],
     ['94Manage 學員管理', '基本課程管理', '繳費記錄', '成績管理', 'Email 支援', '—', '—', '—']),
    ('標準版', 'NT$5,999/月', '200 學員以下', COLORS['primary'],
     ['全部入門版功能', '94inClass 點名系統', 'NFC + 手動點名', 'LINE 家長通知', '出勤報表', 'AI 流失預警', '優先支援', '—']),
    ('專業版', 'NT$8,999/月', '500 學員以下', COLORS['accent'],
     ['全部標準版功能', '94Stock 庫存管理', 'AI 臉辨點名', 'Telegram Bot 操作', 'AI 備貨預測', '多分校支援', '專屬客服', 'API 整合']),
]

for i, (name, price, cap, color, features) in enumerate(plans):
    left = Inches(0.5 + i * 4.2)
    top = Inches(1.5)

    # 推薦標記
    if i == 1:
        add_rounded_rect(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.2), Inches(0.35), COLORS['accent'])
        add_text(slide, left + Inches(0.8), top - Inches(0.13), Inches(2.2), Inches(0.33),
                 '⭐ 最受歡迎', font_size=12, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

    card = add_rounded_rect(slide, left, top + Inches(0.15), Inches(3.8), Inches(5.7), COLORS['white'])
    # 頂部色帶
    add_rect(slide, left, top + Inches(0.15), Inches(3.8), Inches(1.3), color)
    add_text(slide, left, top + Inches(0.3), Inches(3.8), Inches(0.4),
             name, font_size=22, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.75), Inches(3.8), Inches(0.4),
             price, font_size=28, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.15), Inches(3.8), Inches(0.3),
             cap, font_size=13, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

    # 功能列表
    for j, feat in enumerate(features):
        ft = Inches(1.7) + top + Inches(j * 0.45)
        if feat == '—':
            add_text(slide, left + Inches(0.4), ft, Inches(3.0), Inches(0.35),
                     '—', font_size=13, color=COLORS['light2'])
        else:
            add_text(slide, left + Inches(0.4), ft, Inches(3.0), Inches(0.35),
                     f'✓  {feat}', font_size=13, color=COLORS['green_check'])


# =========================================================
# SLIDE 14: 導入流程
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['light'])
slide_header(slide, '輕鬆導入流程', '4 步驟，最快當天開始使用')

steps = [
    ('1', '免費諮詢', '了解您的需求\n推薦適合方案', COLORS['primary']),
    ('2', '帳號開通', '30 分鐘完成設定\n匯入現有學員資料', COLORS['secondary']),
    ('3', '教育訓練', '1 小時快速上手\n提供操作手冊', COLORS['accent']),
    ('4', '正式啟用', '30 天免費試用\n隨時都有支援', COLORS['accent2']),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * 3.3)
    top = Inches(2.2)

    # 圓形步驟編號
    circle = add_circle(slide, left + Inches(1.05), top, Inches(1.0), color)
    add_text(slide, left + Inches(1.05), top + Inches(0.1), Inches(1.0), Inches(0.8),
             num, font_size=36, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

    # 連接線
    if i < 3:
        add_rect(slide, left + Inches(2.15), top + Inches(0.45), Inches(1.3), Inches(0.04), color)

    # 標題
    add_text(slide, left, top + Inches(1.2), Inches(3.1), Inches(0.4),
             title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text(slide, left + Inches(0.2), top + Inches(1.7), Inches(2.7), Inches(1.0),
             desc, font_size=14, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)

# 底部承諾
promises = [
    ('✓ 免費資料轉移', '舊系統資料完整匯入'),
    ('✓ 免費教育訓練', '保證所有員工會操作'),
    ('✓ 30 天無條件退費', '不滿意隨時取消'),
    ('✓ 24hr 技術支援', '問題隨時幫您解決'),
]

for i, (title, desc) in enumerate(promises):
    left = Inches(0.5 + i * 3.2)
    top = Inches(5.3)
    card = add_rounded_rect(slide, left, top, Inches(2.9), Inches(1.2), COLORS['white'])
    add_text(slide, left, top + Inches(0.15), Inches(2.9), Inches(0.35),
             title, font_size=15, color=COLORS['green_check'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.55), Inches(2.9), Inches(0.35),
             desc, font_size=12, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 15: CTA 結尾
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS['dark'])

# 裝飾
add_circle(slide, Inches(-2), Inches(-2), Inches(7), COLORS['dark2'])
add_circle(slide, Inches(10), Inches(4), Inches(6), COLORS['dark2'])
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), COLORS['accent'])

# 主文
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
         '準備好升級您的補習班了嗎？', font_size=20, color=COLORS['light2'])

add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
         '讓 94Cram 成為您最強大的經營武器', font_size=46, color=COLORS['white'], bold=True)

add_rect(slide, Inches(1), Inches(3.5), Inches(3), Inches(0.04), COLORS['accent'])

# 聯絡資訊
contact_info = [
    ('🌐', '94cram.com'),
    ('📧', 'contact@94cram.com'),
    ('📱', 'LINE 官方帳號 @94cram'),
    ('💬', 'Telegram @bot94cram'),
]

for i, (icon, info) in enumerate(contact_info):
    top = Inches(4.0 + i * 0.55)
    add_text(slide, Inches(1.2), top, Inches(0.4), Inches(0.4),
             icon, font_size=18)
    add_text(slide, Inches(1.8), top + Inches(0.02), Inches(5), Inches(0.4),
             info, font_size=18, color=COLORS['light2'])

# CTA 按鈕
add_rounded_rect(slide, Inches(8.0), Inches(3.8), Inches(4.5), Inches(1.5), COLORS['accent'])
add_text(slide, Inches(8.0), Inches(4.0), Inches(4.5), Inches(0.5),
         '立即免費試用', font_size=28, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.0), Inches(4.5), Inches(4.5), Inches(0.4),
         '30 天完整體驗 · 零風險', font_size=16, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(8.0), Inches(5.5), Inches(4.5), Inches(1.0), COLORS['primary'])
add_text(slide, Inches(8.0), Inches(5.6), Inches(4.5), Inches(0.5),
         '預約 Demo 演示', font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.0), Inches(6.05), Inches(4.5), Inches(0.35),
         '專人為您展示完整功能', font_size=14, color=COLORS['light2'], alignment=PP_ALIGN.CENTER)

# 底部
add_text(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.4),
         '© 2026 94Cram 智慧補教管理系統  |  三大系統 × AI 驅動 × 一站式解決方案',
         font_size=12, color=COLORS['text_light'], alignment=PP_ALIGN.CENTER)


# =========================================================
# 儲存檔案
# =========================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '94Cram_行銷簡報_Demo.pptx')
prs.save(output_path)
print(f'✅ 簡報已生成：{output_path}')
print(f'📊 共 {len(prs.slides)} 頁投影片')
